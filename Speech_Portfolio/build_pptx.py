"""
Builds the 8-slide PowerPoint deck for the Independent Speech project.
Theme: Navy (#0B2545) + Gold (#C9A227) + Off-white (#F5F7FA) + Charcoal (#1F2937)
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.oxml.ns import qn
from lxml import etree

# ---------- THEME ----------
NAVY = RGBColor(0x0B, 0x25, 0x45)
GOLD = RGBColor(0xC9, 0xA2, 0x27)
OFFWHITE = RGBColor(0xF5, 0xF7, 0xFA)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
CHARCOAL = RGBColor(0x1F, 0x29, 0x37)
GREY = RGBColor(0x6B, 0x72, 0x80)
LIGHT_GREY = RGBColor(0xD1, 0xD5, 0xDB)

HEADING_FONT = "Georgia"
BODY_FONT = "Calibri"


# ---------- HELPERS ----------
def set_bg(slide, color):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_rect(slide, left, top, width, height, fill_color, line_color=None):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if line_color:
        shape.line.color.rgb = line_color
        shape.line.width = Pt(0.5)
    else:
        shape.line.fill.background()
    shape.shadow.inherit = False
    return shape


def add_text(slide, left, top, width, height, text, *,
             font=BODY_FONT, size=18, bold=False, italic=False,
             color=CHARCOAL, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP):
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0)
    tf.margin_right = Inches(0)
    tf.margin_top = Inches(0)
    tf.margin_bottom = Inches(0)
    tf.vertical_anchor = anchor

    lines = text.split("\n")
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        run = p.add_run()
        run.text = line
        run.font.name = font
        run.font.size = Pt(size)
        run.font.bold = bold
        run.font.italic = italic
        run.font.color.rgb = color
    return tb


def add_speaker_notes(slide, text):
    notes = slide.notes_slide.notes_text_frame
    notes.text = text


def slide_dims(prs):
    return prs.slide_width, prs.slide_height


def add_gold_rule(slide, left, top, width, thickness_pt=2):
    height = Pt(thickness_pt)
    return add_rect(slide, left, top, width, height, GOLD)


def add_footer(slide, prs, slide_num, total=8):
    sw, sh = slide_dims(prs)
    add_text(slide, Inches(0.4), sh - Inches(0.4),
             Inches(6), Inches(0.3),
             "Independent Speech & Public Speaking Portfolio",
             size=9, color=GREY, italic=True)
    add_text(slide, sw - Inches(1.5), sh - Inches(0.4),
             Inches(1.1), Inches(0.3),
             f"{slide_num} / {total}",
             size=9, color=GREY, align=PP_ALIGN.RIGHT)


# ---------- BUILD ----------
prs = Presentation()
prs.slide_width = Inches(13.333)   # 16:9 widescreen
prs.slide_height = Inches(7.5)
SW = prs.slide_width
SH = prs.slide_height

blank_layout = prs.slide_layouts[6]


# ============= SLIDE 1 — TITLE =============
s1 = prs.slides.add_slide(blank_layout)
set_bg(s1, NAVY)

# Eyebrow (gold)
add_text(s1, Inches(0), Inches(1.6), SW, Inches(0.4),
         "BACHELOR OF FINANCIAL MARKETS  |  INDEPENDENT SPEECH",
         font=BODY_FONT, size=14, bold=True, color=GOLD, align=PP_ALIGN.CENTER)

# Title
add_text(s1, Inches(0.5), Inches(2.4), Inches(12.3), Inches(2.0),
         "Why Financial Literacy",
         font=HEADING_FONT, size=54, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
add_text(s1, Inches(0.5), Inches(3.3), Inches(12.3), Inches(2.0),
         "Is the Most Important Skill for My Generation",
         font=HEADING_FONT, size=36, bold=False, color=WHITE, align=PP_ALIGN.CENTER)

# Gold rule
add_gold_rule(s1, Inches(5.67), Inches(4.7), Inches(2.0), 2)

# Subtitle
add_text(s1, Inches(0.5), Inches(5.0), Inches(12.3), Inches(0.5),
         "An Independent Speech by [Student Name]  |  BFM",
         font=BODY_FONT, size=18, italic=True, color=OFFWHITE, align=PP_ALIGN.CENTER)

add_speaker_notes(s1, "Pause for 2 seconds. Smile. Then begin with the hook question.")


# ============= SLIDE 2 — THE HOOK =============
s2 = prs.slides.add_slide(blank_layout)
set_bg(s2, OFFWHITE)

# Section label
add_text(s2, Inches(0.6), Inches(0.5), Inches(6), Inches(0.4),
         "01  /  THE HOOK",
         font=BODY_FONT, size=12, bold=True, color=GOLD)
add_gold_rule(s2, Inches(0.6), Inches(0.95), Inches(0.8), 2)

# Big questions
add_text(s2, Inches(0.6), Inches(2.0), Inches(12.1), Inches(1.5),
         "How much money did you",
         font=HEADING_FONT, size=44, bold=True, color=NAVY, align=PP_ALIGN.LEFT)
add_text(s2, Inches(0.6), Inches(2.85), Inches(12.1), Inches(1.5),
         "spend last month?",
         font=HEADING_FONT, size=44, bold=True, color=NAVY, align=PP_ALIGN.LEFT)

# Second question (gold accent)
add_text(s2, Inches(0.6), Inches(4.4), Inches(12.1), Inches(1.5),
         "And — where did it actually go?",
         font=HEADING_FONT, size=32, italic=True, color=GOLD, align=PP_ALIGN.LEFT)

add_footer(s2, prs, 2)
add_speaker_notes(s2, "Let the questions sit. Don't rush. Let the audience think.")


# ============= SLIDE 3 — THE PROBLEM =============
s3 = prs.slides.add_slide(blank_layout)
set_bg(s3, OFFWHITE)

add_text(s3, Inches(0.6), Inches(0.5), Inches(6), Inches(0.4),
         "02  /  THE PROBLEM",
         size=12, bold=True, color=GOLD)
add_gold_rule(s3, Inches(0.6), Inches(0.95), Inches(0.8), 2)

# Left side - stat
add_text(s3, Inches(0.6), Inches(1.6), Inches(7), Inches(1.2),
         "Only",
         font=HEADING_FONT, size=32, color=CHARCOAL)
add_text(s3, Inches(0.6), Inches(2.0), Inches(7), Inches(2.5),
         "27%",
         font=HEADING_FONT, size=140, bold=True, color=NAVY)
add_text(s3, Inches(0.6), Inches(4.6), Inches(7), Inches(0.6),
         "of Indians are financially literate.",
         font=HEADING_FONT, size=22, color=CHARCOAL)
add_text(s3, Inches(0.6), Inches(5.3), Inches(7), Inches(1.5),
         "That means 3 out of every 4 people\naround you don't fully understand\ntheir own money.",
         size=16, italic=True, color=GREY)

# Right side - 4 person icons
icon_top = Inches(2.2)
icon_size = Inches(0.9)
gap = Inches(0.25)
start_left = Inches(8.4)
for i in range(4):
    color = GOLD if i == 0 else LIGHT_GREY
    left = start_left + (icon_size + gap) * i
    # Head
    head = s3.shapes.add_shape(MSO_SHAPE.OVAL, left + Inches(0.27), icon_top,
                                Inches(0.36), Inches(0.36))
    head.fill.solid()
    head.fill.fore_color.rgb = color
    head.line.fill.background()
    # Body
    body = s3.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                left + Inches(0.05), icon_top + Inches(0.45),
                                Inches(0.8), Inches(1.1))
    body.fill.solid()
    body.fill.fore_color.rgb = color
    body.line.fill.background()

# Caption
add_text(s3, Inches(8.4), Inches(4.0), Inches(4.5), Inches(0.5),
         "1 in 4 understands money.",
         size=14, bold=True, color=NAVY, align=PP_ALIGN.LEFT)
add_text(s3, Inches(8.4), Inches(4.4), Inches(4.5), Inches(0.5),
         "3 in 4 don't.",
         size=14, color=GREY, align=PP_ALIGN.LEFT)

add_footer(s3, prs, 3)
add_speaker_notes(s3, "Emphasize the number. Pause before the next slide.")


# ============= SLIDE 4 — WHAT IT REALLY MEANS =============
s4 = prs.slides.add_slide(blank_layout)
set_bg(s4, OFFWHITE)

add_text(s4, Inches(0.6), Inches(0.5), Inches(8), Inches(0.4),
         "03  /  DEFINITION",
         size=12, bold=True, color=GOLD)
add_gold_rule(s4, Inches(0.6), Inches(0.95), Inches(0.8), 2)

add_text(s4, Inches(0.6), Inches(1.4), Inches(12), Inches(1.0),
         "Financial Literacy = A Life Skill",
         font=HEADING_FONT, size=40, bold=True, color=NAVY, align=PP_ALIGN.CENTER)

# Three columns
col_top = Inches(3.2)
col_h = Inches(3.3)
col_w = Inches(3.6)
gap = Inches(0.4)
start = Inches(0.95)

cards = [
    ("UNDERSTAND", "Know what money is, where it comes from, and what it does."),
    ("MANAGE", "Track, budget, and control your spending without stress."),
    ("GROW", "Save, invest, and build wealth — confidently and patiently."),
]

for i, (head, body) in enumerate(cards):
    left = start + (col_w + gap) * i
    # Card
    card = add_rect(s4, left, col_top, col_w, col_h, WHITE, line_color=LIGHT_GREY)
    # Number circle
    circle = s4.shapes.add_shape(MSO_SHAPE.OVAL, left + Inches(0.4), col_top + Inches(0.4),
                                  Inches(0.6), Inches(0.6))
    circle.fill.solid()
    circle.fill.fore_color.rgb = GOLD
    circle.line.fill.background()
    add_text(s4, left + Inches(0.4), col_top + Inches(0.4), Inches(0.6), Inches(0.6),
             str(i + 1), font=HEADING_FONT, size=22, bold=True, color=WHITE,
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    # Title
    add_text(s4, left + Inches(0.3), col_top + Inches(1.2), col_w - Inches(0.6), Inches(0.6),
             head, font=HEADING_FONT, size=22, bold=True, color=NAVY)
    # Body
    add_text(s4, left + Inches(0.3), col_top + Inches(1.9), col_w - Inches(0.6), Inches(1.3),
             body, size=14, color=CHARCOAL)

add_footer(s4, prs, 4)
add_speaker_notes(s4, "Stress: 'It's not Buffett-level finance. It's daily-life finance.'")


# ============= SLIDE 5 — WHY NOW =============
s5 = prs.slides.add_slide(blank_layout)
set_bg(s5, OFFWHITE)

add_text(s5, Inches(0.6), Inches(0.5), Inches(8), Inches(0.4),
         "04  /  WHY NOW",
         size=12, bold=True, color=GOLD)
add_gold_rule(s5, Inches(0.6), Inches(0.95), Inches(0.8), 2)

add_text(s5, Inches(0.6), Inches(1.5), Inches(12), Inches(1.0),
         "Our money moves faster than our understanding.",
         font=HEADING_FONT, size=32, bold=True, color=NAVY, align=PP_ALIGN.LEFT)

# Bullet list
bullets = [
    ("UPI, instant loans, app-based investing", "Money moves in 2 seconds. Decisions don't catch up."),
    ("Scams, impulsive spending, inflation", "Same phone that helps you save can also empty your wallet."),
    ("Time is our biggest advantage", "Compounding rewards the early — not the big."),
]
y = Inches(3.0)
for title, sub in bullets:
    # Gold dot
    dot = s5.shapes.add_shape(MSO_SHAPE.OVAL, Inches(0.7), y + Inches(0.15),
                               Inches(0.18), Inches(0.18))
    dot.fill.solid()
    dot.fill.fore_color.rgb = GOLD
    dot.line.fill.background()
    add_text(s5, Inches(1.1), y, Inches(11.5), Inches(0.5),
             title, size=18, bold=True, color=NAVY)
    add_text(s5, Inches(1.1), y + Inches(0.45), Inches(11.5), Inches(0.5),
             sub, size=13, italic=True, color=GREY)
    y += Inches(1.15)

add_footer(s5, prs, 5)
add_speaker_notes(s5, "Use the samosa example here for warmth and humor.")


# ============= SLIDE 6 — POWER OF STARTING EARLY =============
s6 = prs.slides.add_slide(blank_layout)
set_bg(s6, NAVY)

add_text(s6, Inches(0.6), Inches(0.5), Inches(8), Inches(0.4),
         "05  /  COMPOUNDING",
         size=12, bold=True, color=GOLD)
add_gold_rule(s6, Inches(0.6), Inches(0.95), Inches(0.8), 2)

add_text(s6, Inches(0.6), Inches(1.4), Inches(12.1), Inches(1.2),
         "Compounding rewards the early,\nnot the big.",
         font=HEADING_FONT, size=36, bold=True, color=WHITE, align=PP_ALIGN.LEFT)

add_text(s6, Inches(0.6), Inches(3.0), Inches(12.1), Inches(0.6),
         "₹2,000/month from age 19 → can grow into something life-changing by 45.",
         size=18, italic=True, color=OFFWHITE, align=PP_ALIGN.LEFT)

# Simple ascending chart - 6 bars
chart_left = Inches(0.7)
chart_top = Inches(4.0)
chart_w = Inches(12.0)
chart_h = Inches(2.7)

# Baseline
baseline = chart_top + chart_h
add_rect(s6, chart_left, baseline, chart_w, Pt(1.5), GOLD)

# Bars
heights_pct = [0.10, 0.18, 0.30, 0.45, 0.65, 0.95]
n = len(heights_pct)
bar_w = Inches(1.1)
gap = (chart_w - bar_w * n) / (n - 1)
ages = ["19", "24", "29", "34", "39", "45"]
for i, h in enumerate(heights_pct):
    bh = chart_h * h
    left = chart_left + (bar_w + gap) * i
    top = baseline - bh
    color = GOLD if i == n - 1 else WHITE
    bar = s6.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, bar_w, bh)
    bar.fill.solid()
    bar.fill.fore_color.rgb = color
    bar.line.fill.background()
    # Age label
    add_text(s6, left, baseline + Inches(0.1), bar_w, Inches(0.4),
             f"Age {ages[i]}", size=11, color=OFFWHITE, align=PP_ALIGN.CENTER)

add_footer(s6, prs, 6)
add_speaker_notes(s6, "Don't get technical. Say: 'Start early — that's the entire secret.'")


# ============= SLIDE 7 — THREE HABITS =============
s7 = prs.slides.add_slide(blank_layout)
set_bg(s7, OFFWHITE)

add_text(s7, Inches(0.6), Inches(0.5), Inches(8), Inches(0.4),
         "06  /  KEY TAKEAWAYS",
         size=12, bold=True, color=GOLD)
add_gold_rule(s7, Inches(0.6), Inches(0.95), Inches(0.8), 2)

add_text(s7, Inches(0.6), Inches(1.4), Inches(12.1), Inches(1.0),
         "Three Habits.  One Life Skill.",
         font=HEADING_FONT, size=40, bold=True, color=NAVY, align=PP_ALIGN.CENTER)

# Three big tiles
col_top = Inches(3.2)
col_h = Inches(3.4)
col_w = Inches(3.8)
gap = Inches(0.35)
start = Inches(0.7)

habits = [
    ("TRACK", "Every rupee\nfor 30 days."),
    ("SAVE", "Follow the\n50-30-20 rule."),
    ("INVEST", "Start with\na ₹500 SIP."),
]

for i, (head, body) in enumerate(habits):
    left = start + (col_w + gap) * i
    # Card with navy header strip
    card = add_rect(s7, left, col_top, col_w, col_h, WHITE, line_color=LIGHT_GREY)
    add_rect(s7, left, col_top, col_w, Inches(1.1), NAVY)
    # Number
    add_text(s7, left, col_top + Inches(0.18), col_w, Inches(0.4),
             f"0{i + 1}", size=14, bold=True, color=GOLD, align=PP_ALIGN.CENTER)
    # Habit name
    add_text(s7, left, col_top + Inches(0.45), col_w, Inches(0.6),
             head, font=HEADING_FONT, size=26, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    # Body
    add_text(s7, left + Inches(0.3), col_top + Inches(1.4), col_w - Inches(0.6), Inches(2.0),
             body, size=18, color=CHARCOAL, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

add_footer(s7, prs, 7)
add_speaker_notes(s7, "Slow down. Repeat: 'Track. Save. Invest.' Let it land.")


# ============= SLIDE 8 — CLOSE =============
s8 = prs.slides.add_slide(blank_layout)
set_bg(s8, NAVY)

add_text(s8, Inches(0.6), Inches(0.7), Inches(8), Inches(0.4),
         "07  /  THE CLOSE",
         size=12, bold=True, color=GOLD)
add_gold_rule(s8, Inches(0.6), Inches(1.15), Inches(0.8), 2)

add_text(s8, Inches(0.5), Inches(2.0), Inches(12.3), Inches(1.4),
         "It's not about getting rich quickly.",
         font=HEADING_FONT, size=34, italic=True, color=OFFWHITE, align=PP_ALIGN.CENTER)
add_text(s8, Inches(0.5), Inches(2.7), Inches(12.3), Inches(1.4),
         "It's about getting in control quietly.",
         font=HEADING_FONT, size=38, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

add_gold_rule(s8, Inches(5.67), Inches(4.2), Inches(2.0), 2)

add_text(s8, Inches(0.5), Inches(4.5), Inches(12.3), Inches(0.6),
         "Track every rupee you spend for the next 7 days.",
         size=18, italic=True, color=GOLD, align=PP_ALIGN.CENTER)
add_text(s8, Inches(0.5), Inches(5.0), Inches(12.3), Inches(0.6),
         "That's where it begins.",
         size=16, italic=True, color=OFFWHITE, align=PP_ALIGN.CENTER)

add_text(s8, Inches(0.5), Inches(6.0), Inches(12.3), Inches(0.6),
         "THANK YOU",
         font=HEADING_FONT, size=28, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
add_text(s8, Inches(0.5), Inches(6.6), Inches(12.3), Inches(0.4),
         "[Student Name]   |   BFM",
         size=12, color=OFFWHITE, align=PP_ALIGN.CENTER)

add_speaker_notes(s8, "Pause. Maintain eye contact. Smile. Wait for applause before walking off.")


# ---------- SAVE ----------
output = "Hatch/Speech_Portfolio/Independent_Speech_Slide_Deck.pptx"
prs.save(output)
print(f"Saved: {output}")
print(f"Total slides: {len(prs.slides)}")
